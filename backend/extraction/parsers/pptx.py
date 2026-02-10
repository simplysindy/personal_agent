"""PowerPoint parser using python-pptx."""

from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional

from pptx import Presentation
from pptx.util import Inches


@dataclass
class SlideContent:
    """Content extracted from a single slide."""

    slide_number: int
    title: str = ""
    text_content: list[str] = field(default_factory=list)
    speaker_notes: str = ""
    shapes: list[dict] = field(default_factory=list)


@dataclass
class ParsedPptx:
    """Result of parsing a PPTX file."""

    title: str
    text: str
    slide_count: int
    slides: list[SlideContent] = field(default_factory=list)
    core_properties: dict = field(default_factory=dict)


class PptxParser:
    """Parser for PowerPoint presentations using python-pptx."""

    def parse_file(self, file_path: Path) -> ParsedPptx:
        """Parse a PPTX file and extract slides, text, and metadata."""
        try:
            prs = Presentation(file_path)
        except Exception as e:
            raise ValueError(f"Could not open presentation: {file_path} - {e}")

        # Extract core properties
        core_props = {}
        try:
            cp = prs.core_properties
            core_props = {
                "title": cp.title,
                "author": cp.author,
                "subject": cp.subject,
                "keywords": cp.keywords,
                "created": cp.created.isoformat() if cp.created else None,
                "modified": cp.modified.isoformat() if cp.modified else None,
            }
        except Exception:
            pass

        # Extract slides
        slides = []
        all_text = []
        presentation_title = core_props.get("title") or file_path.stem

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = self._extract_slide_content(slide, slide_num)
            slides.append(slide_content)

            # Build text representation
            if slide_content.title:
                all_text.append(f"## Slide {slide_num}: {slide_content.title}")
                # Use first slide title as presentation title if not in metadata
                if slide_num == 1 and not core_props.get("title"):
                    presentation_title = slide_content.title
            else:
                all_text.append(f"## Slide {slide_num}")

            all_text.extend(slide_content.text_content)

            if slide_content.speaker_notes:
                all_text.append(f"[Notes: {slide_content.speaker_notes}]")

        return ParsedPptx(
            title=presentation_title,
            text="\n\n".join(all_text),
            slide_count=len(prs.slides),
            slides=slides,
            core_properties=core_props,
        )

    def _extract_slide_content(self, slide, slide_num: int) -> SlideContent:
        """Extract content from a single slide."""
        title = ""
        text_content = []
        shapes_info = []

        for shape in slide.shapes:
            # Get title
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    title = shape.text.strip()
                else:
                    # Extract all text from shape
                    shape_text = []
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            shape_text.append(para_text)
                    if shape_text:
                        text_content.extend(shape_text)

                shapes_info.append({
                    "type": "text",
                    "content": shape.text.strip(),
                })

            elif shape.has_table:
                # Extract table content
                table = shape.table
                table_text = []
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    table_text.append(row_text)
                if table_text:
                    text_content.append("\n".join(table_text))
                shapes_info.append({
                    "type": "table",
                    "rows": len(table.rows),
                    "cols": len(table.columns),
                })

            elif hasattr(shape, "image"):
                shapes_info.append({
                    "type": "image",
                    "content_type": getattr(shape.image, "content_type", "unknown"),
                })

        # Extract speaker notes
        notes = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes = notes_frame.text.strip()

        return SlideContent(
            slide_number=slide_num,
            title=title,
            text_content=text_content,
            speaker_notes=notes,
            shapes=shapes_info,
        )

    def extract_text_only(self, file_path: Path) -> str:
        """Quick extraction of just the text content."""
        try:
            prs = Presentation(file_path)
            text_parts = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if text:
                            text_parts.append(text)
                    elif shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells
                            )
                            if row_text.strip():
                                text_parts.append(row_text)

                # Include speaker notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame
                    if notes and notes.text.strip():
                        text_parts.append(f"[Notes: {notes.text.strip()}]")

            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Failed to extract text from {file_path}: {e}")
